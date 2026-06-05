from pptx import Presentation
from .base_parser import BaseParser

class PptxParser(BaseParser):
    def parse(self, file_path: str) -> str:
        text = ""
        try:
            prs = Presentation(file_path)
            for i, slide in enumerate(prs.slides):
                text += f"\n--- Slide {i + 1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            print(f"Error parsing PPTX {file_path}: {e}")
        return text
