import json
from pathlib import Path
from typing import Any

from app.core.config import settings


ARTIFACT_VERSION = "1.0"


def artifact_path(document_id: str) -> Path:
    return Path(settings.UPLOAD_DIR) / f"{document_id}.layout.json"


def _text_block(text: str, *, page: int, kind: str = "text", order: int = 0, bbox: list[float] | None = None) -> dict[str, Any]:
    return {"kind": kind, "text": text.strip(), "page": page, "order": order, "bbox": bbox}


def extract_layout(file_path: str) -> dict[str, Any]:
    """Extract a lightweight, provenance-preserving structural representation."""
    ext = Path(file_path).suffix.lower()
    pages: list[dict[str, Any]] = []

    if ext == ".pdf":
        import fitz
        with fitz.open(file_path) as doc:
            for page_no, page in enumerate(doc, start=1):
                blocks = []
                for order, block in enumerate(page.get_text("blocks")):
                    text = str(block[4]).strip()
                    if text:
                        blocks.append(_text_block(text, page=page_no, order=order, bbox=list(block[:4])))
                images = len(page.get_images(full=True))
                pages.append({"page": page_no, "blocks": blocks, "image_count": images})

    elif ext == ".docx":
        from docx import Document as DocxDocument
        doc = DocxDocument(file_path)
        blocks = []
        order = 0
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                kind = "heading" if paragraph.style and paragraph.style.name.lower().startswith("heading") else "text"
                blocks.append(_text_block(paragraph.text, page=1, kind=kind, order=order))
                order += 1
        tables = []
        for table_idx, table in enumerate(doc.tables):
            rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
            tables.append({"index": table_idx, "rows": rows})
        pages.append({"page": 1, "blocks": blocks, "tables": tables, "image_count": len(doc.inline_shapes)})

    elif ext == ".pptx":
        from pptx import Presentation
        presentation = Presentation(file_path)
        for page_no, slide in enumerate(presentation.slides, start=1):
            blocks = []
            order = 0
            image_count = 0
            for shape in slide.shapes:
                if getattr(shape, "shape_type", None) == 13:
                    image_count += 1
                if hasattr(shape, "text") and shape.text.strip():
                    blocks.append(_text_block(shape.text, page=page_no, order=order))
                    order += 1
            pages.append({"page": page_no, "blocks": blocks, "image_count": image_count})

    elif ext in {".xlsx", ".xls"}:
        from openpyxl import load_workbook
        workbook = load_workbook(file_path, read_only=True, data_only=True)
        for sheet_no, sheet in enumerate(workbook.worksheets, start=1):
            rows = []
            for row in sheet.iter_rows(values_only=True):
                values = ["" if value is None else str(value) for value in row]
                if any(values):
                    rows.append(values)
            pages.append({"page": sheet_no, "sheet": sheet.title, "rows": rows, "blocks": []})
        workbook.close()
    else:
        raise ValueError(f"Unsupported multimodal format: {ext}")

    return {"artifact_version": ARTIFACT_VERSION, "file_type": ext, "page_count": len(pages), "pages": pages}


def build_artifact(document_id: str, file_path: str, *, persist: bool = True) -> dict[str, Any]:
    artifact = extract_layout(file_path)
    if persist:
        path = artifact_path(document_id)
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text(json.dumps(artifact, ensure_ascii=False), encoding="utf-8")
    return artifact


def load_artifact(document_id: str) -> dict[str, Any] | None:
    path = artifact_path(document_id)
    if not path.exists():
        return None
    return json.loads(path.read_text(encoding="utf-8"))


def flatten_artifact(artifact: dict[str, Any]) -> list[dict[str, Any]]:
    rows: list[dict[str, Any]] = []
    for page in artifact.get("pages", []):
        for block in page.get("blocks", []):
            if block.get("text"):
                rows.append({"page": page.get("page"), "kind": block.get("kind", "text"), "text": block["text"]})
        for table in page.get("tables", []):
            rows.append({"page": page.get("page"), "kind": "table", "text": "\n".join(" | ".join(row) for row in table.get("rows", []))})
        for row in page.get("rows", []):
            rows.append({"page": page.get("page"), "kind": "spreadsheet_row", "text": " | ".join(row)})
    return rows
